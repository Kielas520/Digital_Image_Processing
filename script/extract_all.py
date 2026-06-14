import os
import sys
import re
import json
import shutil
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

try:
    import fitz  # pymupdf
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

try:
    from docx import Document
    from docx.oxml.ns import qn
except ImportError:
    Document = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None


OMML_NS = {
    'm': 'http://schemas.openxmlformats.org/officeDocument/2006/math',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def omml_to_latex(omml_element):
    latex_parts = []
    _parse_omml_node(omml_element, latex_parts)
    result = ''.join(latex_parts).strip()
    if result:
        return f"${result}$"
    return ""


def _parse_omml_node(node, parts):
    tag = node.tag.split('}')[-1] if '}' in node.tag else node.tag

    if tag == 'r':
        text = ''
        for child in node:
            ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if ctag == 't':
                text += (child.text or '')
        parts.append(text)
    elif tag == 'f':
        num = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}num')
        den = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}den')
        parts.append('\\frac{')
        if num is not None:
            _parse_omml_node(num, parts)
        parts.append('}{')
        if den is not None:
            _parse_omml_node(den, parts)
        parts.append('}')
    elif tag == 'rad':
        deg = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}deg')
        e = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}e')
        deg_text = []
        if deg is not None:
            _parse_omml_node(deg, deg_text)
        deg_str = ''.join(deg_text).strip()
        if deg_str and deg_str != '2':
            parts.append(f'\\sqrt[{deg_str}]{{')
        else:
            parts.append('\\sqrt{')
        if e is not None:
            _parse_omml_node(e, parts)
        parts.append('}')
    elif tag == 'sSup':
        base = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}e')
        sup = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}sup')
        if base is not None:
            _parse_omml_node(base, parts)
        parts.append('^{')
        if sup is not None:
            _parse_omml_node(sup, parts)
        parts.append('}')
    elif tag == 'sSub':
        base = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}e')
        sub = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}sub')
        if base is not None:
            _parse_omml_node(base, parts)
        parts.append('_{')
        if sub is not None:
            _parse_omml_node(sub, parts)
        parts.append('}')
    elif tag == 'nary':
        chr_elem = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}naryPr/{http://schemas.openxmlformats.org/officeDocument/2006/math}chr')
        char = chr_elem.get('{http://schemas.openxmlformats.org/officeDocument/2006/math}val', '∑') if chr_elem is not None else '∑'
        char_map = {'∑': '\\sum', '∏': '\\prod', '∫': '\\int'}
        parts.append(char_map.get(char, char))
        sub = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}sub')
        sup = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}sup')
        e = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}e')
        if sub is not None:
            parts.append('_{')
            _parse_omml_node(sub, parts)
            parts.append('}')
        if sup is not None:
            parts.append('^{')
            _parse_omml_node(sup, parts)
            parts.append('}')
        if e is not None:
            parts.append(' ')
            _parse_omml_node(e, parts)
    elif tag == 'd':
        dPr = node.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}dPr')
        begChr = '('
        endChr = ')'
        if dPr is not None:
            bc = dPr.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}begChr')
            ec = dPr.find('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}endChr')
            if bc is not None:
                begChr = bc.get('{http://schemas.openxmlformats.org/officeDocument/2006/math}val', '(')
            if ec is not None:
                endChr = ec.get('{http://schemas.openxmlformats.org/officeDocument/2006/math}val', ')')
        parts.append(f'\\left{begChr}')
        for child in node:
            ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if ctag == 'e':
                _parse_omml_node(child, parts)
        parts.append(f'\\right{endChr}')
    elif tag == 'm':
        for child in node:
            ctag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if ctag == 'mr':
                row_parts = []
                for mc in child:
                    cell_parts = []
                    _parse_omml_node(mc, cell_parts)
                    row_parts.append(''.join(cell_parts))
                parts.append(' & '.join(row_parts) + ' \\\\ ')
    else:
        for child in node:
            _parse_omml_node(child, parts)


def extract_pptx(filepath, output_dir):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    md_lines = [f"# {Path(filepath).stem}\n"]

    if Presentation is None:
        return _extract_pptx_xml(filepath, output_dir, md_lines)

    try:
        prs = Presentation(filepath)
    except Exception:
        return _extract_pptx_xml(filepath, output_dir, md_lines)

    img_count = 0

    try:
        for slide_idx, slide in enumerate(prs.slides, 1):
            md_lines.append(f"\n## Slide {slide_idx}\n")

            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                md_lines.append(text)
                        md_lines.append("")
                except Exception:
                    pass

                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        header = "| " + " | ".join(cell.text.strip() for cell in table.rows[0].cells) + " |"
                        sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
                        md_lines.append(header)
                        md_lines.append(sep)
                        for row in table.rows[1:]:
                            row_text = "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
                            md_lines.append(row_text)
                        md_lines.append("")
                except Exception:
                    pass

                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_count += 1
                        image = shape.image
                        ext = image.content_type.split('/')[-1]
                        if ext == 'jpeg':
                            ext = 'jpg'
                        img_name = f"slide{slide_idx}_img{img_count}.{ext}"
                        img_path = os.path.join(images_dir, img_name)
                        with open(img_path, 'wb') as f:
                            f.write(image.blob)
                        md_lines.append(f"![{img_name}](images/{img_name})\n")
                except Exception:
                    pass
    except Exception:
        md_lines = [f"# {Path(filepath).stem}\n"]
        return _extract_pptx_xml(filepath, output_dir, md_lines)

    with zipfile.ZipFile(filepath, 'r') as z:
        for name in z.namelist():
            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                xml_content = z.read(name)
                root = ET.fromstring(xml_content)
                slide_num = int(os.path.basename(name).replace('slide', '').replace('.xml', ''))
                math_elements = root.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}oMath')
                if math_elements:
                    for me in math_elements:
                        latex = omml_to_latex(me)
                        if latex:
                            marker = f"## Slide {slide_num}"
                            for i, line in enumerate(md_lines):
                                if line.strip() == marker:
                                    insert_idx = i + 1
                                    while insert_idx < len(md_lines) and md_lines[insert_idx].strip():
                                        insert_idx += 1
                                    md_lines.insert(insert_idx, f"\n{latex}\n")
                                    break

    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[PPTX] {filepath} -> {md_path}")
    return md_path


def _extract_pptx_xml(filepath, output_dir, md_lines):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    img_count = 0

    with zipfile.ZipFile(filepath, 'r') as z:
        for name in z.namelist():
            if name.startswith('ppt/media/'):
                img_count += 1
                filename = os.path.basename(name)
                with z.open(name) as src, open(os.path.join(images_dir, filename), 'wb') as dst:
                    shutil.copyfileobj(src, dst)
                md_lines.append(f"![{filename}](images/{filename})")

        slide_files = sorted(
            [n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')],
            key=lambda x: int(os.path.basename(x).replace('slide', '').replace('.xml', ''))
        )
        for slide_file in slide_files:
            xml_content = z.read(slide_file)
            root = ET.fromstring(xml_content)
            slide_num = os.path.basename(slide_file).replace('slide', '').replace('.xml', '')
            md_lines.append(f"\n## Slide {slide_num}\n")

            for elem in root.iter():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag == 't' and elem.text:
                    md_lines.append(elem.text)

            math_elements = root.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}oMath')
            for me in math_elements:
                latex = omml_to_latex(me)
                if latex:
                    md_lines.append(latex)

    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[PPTX/XML] {filepath} -> {md_path}")
    return md_path


def extract_docx(filepath, output_dir):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    md_lines = [f"# {Path(filepath).stem}\n"]

    if Document is None:
        md_lines.append("*[python-docx not installed, falling back to XML extraction]*\n")
        return _extract_docx_xml(filepath, output_dir, md_lines)

    doc = Document(filepath)
    img_count = 0

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img_count += 1
            ext = rel.target_ref.split('.')[-1]
            img_name = f"img{img_count}.{ext}"
            img_path = os.path.join(images_dir, img_name)
            with open(img_path, 'wb') as f:
                f.write(rel.target_part.blob)
            md_lines.append(f"![{img_name}](images/{img_name})\n")

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style.name.startswith('Heading'):
                level = para.style.name.replace('Heading', '').strip()
                try:
                    level = int(level) + 1
                except ValueError:
                    level = 2
                md_lines.append(f"{'#' * level} {text}\n")
            else:
                md_lines.append(text)
                md_lines.append("")

    for table in doc.tables:
        if table.rows:
            header = "| " + " | ".join(cell.text.strip().replace('\n', ' ') for cell in table.rows[0].cells) + " |"
            sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
            md_lines.append(header)
            md_lines.append(sep)
            for row in table.rows[1:]:
                row_text = "| " + " | ".join(cell.text.strip().replace('\n', ' ') for cell in row.cells) + " |"
                md_lines.append(row_text)
            md_lines.append("")

    with zipfile.ZipFile(filepath, 'r') as z:
        if 'word/document.xml' in z.namelist():
            xml_content = z.read('word/document.xml')
            root = ET.fromstring(xml_content)
            math_elements = root.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}oMath')
            for me in math_elements:
                latex = omml_to_latex(me)
                if latex:
                    md_lines.append(latex)
                    md_lines.append("")

    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[DOCX] {filepath} -> {md_path}")
    return md_path


def _extract_docx_xml(filepath, output_dir, md_lines):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)

    with zipfile.ZipFile(filepath, 'r') as z:
        for name in z.namelist():
            if name.startswith('word/media/'):
                filename = os.path.basename(name)
                with z.open(name) as src, open(os.path.join(images_dir, filename), 'wb') as dst:
                    shutil.copyfileobj(src, dst)
                md_lines.append(f"![{filename}](images/{filename})")

        if 'word/document.xml' in z.namelist():
            xml_content = z.read('word/document.xml')
            root = ET.fromstring(xml_content)
            for elem in root.iter():
                tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag == 't' and elem.text:
                    md_lines.append(elem.text)

            math_elements = root.findall('.//{http://schemas.openxmlformats.org/officeDocument/2006/math}oMath')
            for me in math_elements:
                latex = omml_to_latex(me)
                if latex:
                    md_lines.append(latex)

    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[DOCX/XML] {filepath} -> {md_path}")
    return md_path


def extract_pdf(filepath, output_dir):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    md_lines = [f"# {Path(filepath).stem}\n"]

    if fitz is None:
        md_lines.append("*[pymupdf not installed, cannot extract PDF]*\n")
        md_path = os.path.join(output_dir, 'content.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(md_lines))
        print(f"[PDF] SKIPPED (no pymupdf): {filepath}")
        return md_path

    doc = fitz.open(filepath)
    img_count = 0

    for page_idx, page in enumerate(doc, 1):
        md_lines.append(f"\n## Page {page_idx}\n")

        text = page.get_text("text")
        if text.strip():
            md_lines.append(text.strip())
            md_lines.append("")

        images = page.get_images(full=True)
        for img_info in images:
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                img_count += 1
                ext = base_image["ext"]
                img_name = f"page{page_idx}_img{img_count}.{ext}"
                img_path = os.path.join(images_dir, img_name)
                with open(img_path, 'wb') as f:
                    f.write(base_image["image"])
                md_lines.append(f"![{img_name}](images/{img_name})\n")
            except Exception:
                pass

    doc.close()
    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[PDF] {filepath} -> {md_path}")
    return md_path


def extract_html(filepath, output_dir):
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)
    md_lines = [f"# {Path(filepath).stem}\n"]

    if BeautifulSoup is None:
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL)
        content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL)
        content = re.sub(r'<[^>]+>', '\n', content)
        content = re.sub(r'\n{3,}', '\n\n', content)
        md_lines.append(content.strip())
    else:
        with open(filepath, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        for style in soup.find_all(['style', 'script']):
            style.decompose()

        for math_elem in soup.find_all(class_=re.compile(r'katex|math|MathJax')):
            annotation = math_elem.find('annotation', encoding='application/x-tex')
            if annotation and annotation.string:
                math_elem.replace_with(f" ${annotation.string.strip()}$ ")
                continue
            semantics = math_elem.find('semantics')
            if semantics:
                ann = semantics.find('annotation')
                if ann and ann.string:
                    math_elem.replace_with(f" ${ann.string.strip()}$ ")
                    continue
            math_text = math_elem.get_text()
            if math_text.strip():
                math_elem.replace_with(f" ${math_text.strip()}$ ")

        for h_tag in soup.find_all(re.compile(r'^h[1-6]$')):
            level = int(h_tag.name[1])
            text = h_tag.get_text().strip()
            if text:
                md_lines.append(f"\n{'#' * (level + 1)} {text}\n")
            h_tag.decompose()

        for code_block in soup.find_all('pre'):
            code = code_block.find('code')
            if code:
                lang = ''
                if code.get('class'):
                    for cls in code['class']:
                        if cls.startswith('language-'):
                            lang = cls.replace('language-', '')
                md_lines.append(f"\n```{lang}")
                md_lines.append(code.get_text())
                md_lines.append("```\n")
            code_block.decompose()

        for table in soup.find_all('table'):
            rows = table.find_all('tr')
            if rows:
                cells = rows[0].find_all(['th', 'td'])
                header = "| " + " | ".join(c.get_text().strip() for c in cells) + " |"
                sep = "| " + " | ".join("---" for _ in cells) + " |"
                md_lines.append(header)
                md_lines.append(sep)
                for row in rows[1:]:
                    cells = row.find_all(['td', 'th'])
                    md_lines.append("| " + " | ".join(c.get_text().strip() for c in cells) + " |")
                md_lines.append("")
            table.decompose()

        for li in soup.find_all('li'):
            text = li.get_text().strip()
            if text:
                md_lines.append(f"- {text}")

        for p in soup.find_all(['p', 'div']):
            text = p.get_text().strip()
            if text:
                md_lines.append(text)
                md_lines.append("")

    md_path = os.path.join(output_dir, 'content.md')
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md_lines))
    print(f"[HTML] {filepath} -> {md_path}")
    return md_path


def process_file(filepath, base_source_dir, base_output_dir):
    rel_path = os.path.relpath(filepath, base_source_dir)
    stem = Path(filepath).stem
    parent_rel = os.path.dirname(rel_path)
    output_dir = os.path.join(base_output_dir, parent_rel, stem)
    os.makedirs(output_dir, exist_ok=True)

    ext = Path(filepath).suffix.lower()
    try:
        if ext == '.pptx':
            return extract_pptx(filepath, output_dir)
        elif ext == '.docx':
            return extract_docx(filepath, output_dir)
        elif ext == '.pdf':
            return extract_pdf(filepath, output_dir)
        elif ext in ('.html', '.htm'):
            return extract_html(filepath, output_dir)
        else:
            print(f"[SKIP] Unsupported format: {filepath}")
            return None
    except Exception as e:
        print(f"[ERROR] {filepath}: {e}")
        return None


def batch_extract(source_dir, output_dir):
    supported_exts = {'.pptx', '.docx', '.pdf', '.html', '.htm'}
    files = []
    for root, dirs, filenames in os.walk(source_dir):
        for fname in filenames:
            if Path(fname).suffix.lower() in supported_exts:
                files.append(os.path.join(root, fname))

    files.sort()
    print(f"Found {len(files)} files to process.\n")

    results = []
    for f in files:
        result = process_file(f, source_dir, output_dir)
        if result:
            results.append(result)

    print(f"\nDone. Processed {len(results)}/{len(files)} files.")
    print(f"Output directory: {output_dir}")
    return results


if __name__ == "__main__":
    project_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    source_dir = os.path.join(project_dir, 'source')
    output_dir = os.path.join(project_dir, 'source', 'extracted')

    if len(sys.argv) > 1:
        if sys.argv[1] == '--single' and len(sys.argv) > 2:
            filepath = sys.argv[2]
            out = sys.argv[3] if len(sys.argv) > 3 else output_dir
            process_file(filepath, source_dir, out)
        else:
            source_dir = sys.argv[1]
            output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.join(source_dir, 'extracted')
            batch_extract(source_dir, output_dir)
    else:
        batch_extract(source_dir, output_dir)
